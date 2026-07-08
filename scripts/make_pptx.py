"""Generate PITCH_SLIDES.pptx from the pitch content (slide theme).

Run: python scripts/make_pptx.py
Output: PITCH_SLIDES.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Design tokens (match the web UI)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)
INK = RGBColor(0x0A, 0x0A, 0x0A)
INK_SOFT = RGBColor(0x44, 0x40, 0x3C)
MUTED = RGBColor(0x78, 0x71, 0x6C)
ACCENT = RGBColor(0xDC, 0x26, 0x26)
RULE_SOFT = RGBColor(0xD6, 0xD3, 0xD1)

FONT_DISPLAY = "Space Grotesk"
FONT_MONO = "IBM Plex Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, font=FONT_DISPLAY, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_mono_label(slide, x, y, text, color=ACCENT, size=10):
    return add_text(slide, x, y, Inches(4), Inches(0.3), text,
                    font=FONT_MONO, size=size, color=color, bold=True)


def slide_base(paper=True):
    s = prs.slides.add_slide(blank)
    bg = add_rect(s, 0, 0, SW, SH, PAPER if paper else INK)
    return s


def section_header(slide, num, title):
    add_mono_label(slide, Inches(0.8), Inches(0.5), num, color=ACCENT)
    add_text(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.8),
             title, font=FONT_DISPLAY, size=36, bold=True, color=INK)
    # accent underline
    add_rect(slide, Inches(0.8), Inches(1.5), Inches(2), Pt(3), ACCENT)


def page_number(slide, n, total):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3),
             f"{n:02d} / {total:02d}", font=FONT_MONO, size=9, color=MUTED,
             align=PP_ALIGN.RIGHT)


def footer(slide):
    add_text(slide, Inches(0.8), Inches(7.05), Inches(6), Inches(0.3),
             "CarSpec AI · Module 2 Pitch", font=FONT_MONO, size=9, color=MUTED)


TOTAL = 10

# Slide 1: Title
s = slide_base(paper=False)
add_rect(s, 0, 0, SW, SH, INK)
# accent bar
add_rect(s, Inches(0.8), Inches(1.0), Inches(0.5), Pt(3), ACCENT)
add_mono_label(s, Inches(0.8), Inches(1.2), "MODULE 2 · COMPUTER VISION", color=ACCENT)
add_text(s, Inches(0.8), Inches(1.6), Inches(11), Inches(2.5),
         "CarSpec AI", font=FONT_DISPLAY, size=96, bold=True, color=PAPER)
add_text(s, Inches(0.8), Inches(3.8), Inches(11), Inches(1.2),
         "Read a car like a spec sheet.\nFrom one photo.",
         font=FONT_DISPLAY, size=32, color=PAPER, italic=True)
# accent line
add_rect(s, Inches(0.8), Inches(5.6), Inches(3), Pt(2), ACCENT)
add_text(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "5-minute Pitch · 2026", font=FONT_MONO, size=12, color=PAPER)
add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
         "hanfuzhao781-carspec-ai.hf.space", font=FONT_MONO, size=11, color=ACCENT)
page_number(s, 1, TOTAL)

# Slide 2: Problem
s = slide_base()
section_header(s, "01 / PROBLEM", "Problem")
footer(s); page_number(s, 2, TOTAL)

# blockquote
add_rect(s, Inches(0.8), Inches(2.0), Pt(4), Inches(1.2), ACCENT)
add_text(s, Inches(1.1), Inches(2.0), Inches(11), Inches(1.2),
         "Traditional vehicle recognition predicts one attribute and behaves like a black box.",
         font=FONT_DISPLAY, size=24, italic=True, color=INK_SOFT)

# two pains
add_text(s, Inches(0.8), Inches(3.6), Inches(12), Inches(0.4),
         "TWO PAINS", font=FONT_MONO, size=11, color=ACCENT, bold=True)

pains = [
    ("01", "Single-task", "Type or brand or color - never together. Forces separate models per attribute."),
    ("02", "Black box", "Users cannot see why a call was made. No audit trail for predictions."),
]
for i, (n, title, desc) in enumerate(pains):
    x = Inches(0.8 + i * 6.2)
    add_rect(s, x, Inches(4.1), Inches(5.8), Inches(2.0), PAPER, line=INK)
    add_rect(s, x, Inches(4.1), Inches(5.8), Pt(3), INK)
    add_text(s, x + Inches(0.3), Inches(4.3), Inches(0.6), Inches(0.4),
             n, font=FONT_MONO, size=14, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.3), Inches(4.65), Inches(5.2), Inches(0.5),
             title, font=FONT_DISPLAY, size=22, bold=True)
    add_text(s, x + Inches(0.3), Inches(5.2), Inches(5.2), Inches(0.8),
             desc, font=FONT_DISPLAY, size=13, color=INK_SOFT)

# Slide 3: Hypothesis
s = slide_base()
section_header(s, "02 / HYPOTHESIS", "Hypothesis")
footer(s); page_number(s, 3, TOTAL)

add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(0.6),
         "Vehicle attributes are correlated - exploit them via multi-task learning.",
         font=FONT_DISPLAY, size=22, color=INK_SOFT)

# correlation table
rows = [("Type", "Doors", "Seats"),
        ("coupe", "2", "2"),
        ("sedan", "4", "5"),
        ("MPV", "5", "7")]
col_x = [Inches(0.8), Inches(5.0), Inches(9.2)]
col_w = [Inches(4.0), Inches(4.0), Inches(3.4)]
y = Inches(3.0)
rh = Inches(0.7)
for ri, row in enumerate(rows):
    is_head = ri == 0
    fill = INK if is_head else (PAPER if ri % 2 == 1 else RGBColor(0xF0, 0xED, 0xE7))
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], y, col_w[ci], rh, fill, line=RULE_SOFT)
        add_text(s, col_x[ci] + Inches(0.2), y, col_w[ci] - Inches(0.4), rh,
                 cell, font=FONT_MONO if is_head else FONT_DISPLAY,
                 size=14 if is_head else 18, bold=is_head,
                 color=PAPER if is_head else INK, anchor=MSO_ANCHOR.MIDDLE)
    y += rh

add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.6),
         "A multi-task model with shared backbone should pick up these correlations and push accuracy up - and the handcrafted features leave a trail a human can audit.",
         font=FONT_DISPLAY, size=14, italic=True, color=INK_SOFT)

# Slide 4: Method - Three Models
s = slide_base()
section_header(s, "03 / METHOD", "Three Models")
footer(s); page_number(s, 4, TOTAL)

headers = ["", "Naive", "Classical", "Deep"]
data = [
    ("Algorithm", "Majority class", "Random Forest", "MobileNetV2"),
    ("Input", "-", "50-D handcrafted", "224×224 RGB"),
    ("Heads", "1", "3 (independent)", "3 (shared)"),
    ("Why", "Sanity floor", "Interpretable", "Multi-task joint"),
]
col_x = [Inches(0.8), Inches(3.5), Inches(6.7), Inches(10.0)]
col_w = [Inches(2.5), Inches(3.0), Inches(3.0), Inches(2.7)]
y = Inches(2.2)
rh = Inches(0.65)
# header row
for ci, h in enumerate(headers):
    add_rect(s, col_x[ci], y, col_w[ci], rh, INK, line=INK)
    add_text(s, col_x[ci] + Inches(0.2), y, col_w[ci] - Inches(0.4), rh,
             h, font=FONT_MONO, size=12, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)
y += rh
for ri, row in enumerate(data):
    fill = PAPER if ri % 2 == 0 else RGBColor(0xF0, 0xED, 0xE7)
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], y, col_w[ci], rh, fill, line=RULE_SOFT)
        is_label = ci == 0
        add_text(s, col_x[ci] + Inches(0.2), y, col_w[ci] - Inches(0.4), rh,
                 cell, font=FONT_MONO if is_label else FONT_DISPLAY,
                 size=11 if is_label else 13, bold=is_label,
                 color=ACCENT if is_label else INK, anchor=MSO_ANCHOR.MIDDLE)
    y += rh

add_text(s, Inches(0.8), Inches(6.2), Inches(12), Inches(0.6),
         "Backbone: MobileNetV2 (ImageNet pretrained) → 256-D FC → 3 heads (5/3/3 classes).",
         font=FONT_MONO, size=12, color=INK_SOFT)

# Slide 5: Pipeline
s = slide_base()
section_header(s, "04 / PIPELINE", "Pipeline")
footer(s); page_number(s, 5, TOTAL)

steps = [
    ("STEP_01", "CompCars", "136,726 images\nraw input"),
    ("STEP_02", "Preprocess", "224×224 resize\nnormalize"),
    ("STEP_03", "Extract", "50-D handcrafted\nHSV · HOG · LBP"),
    ("STEP_04", "Predict", "RF + MobileNetV2\nmulti-task heads"),
    ("STEP_05", "Explain", "top-5 · confidence\nfeature narrative"),
]
n = len(steps)
margin = 0.8
gap = 0.3
total_w = 13.333 - 2 * margin - (n - 1) * gap
cw = total_w / n
y = Inches(2.5)
ch = Inches(3.0)
for i, (num, title, desc) in enumerate(steps):
    x = Inches(margin + i * (cw + gap))
    add_rect(s, x, y, Inches(cw), ch, PAPER, line=INK)
    add_rect(s, x, y, Inches(cw), Pt(3), ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(cw - 0.4), Inches(0.3),
             num, font=FONT_MONO, size=10, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.6), Inches(cw - 0.4), Inches(0.6),
             title, font=FONT_DISPLAY, size=18, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.4), Inches(cw - 0.4), Inches(1.4),
             desc, font=FONT_DISPLAY, size=11, color=INK_SOFT)
    # arrow
    if i < n - 1:
        ax = x + Inches(cw) + Inches(0.05)
        add_text(s, ax, y + Inches(ch.emu / 914400 / 2 - 0.15), Inches(0.3), Inches(0.3),
                 "→", font=FONT_MONO, size=16, color=ACCENT, bold=True)

add_text(s, Inches(0.8), Inches(6.2), Inches(12), Inches(0.8),
         "50-D features: HSV histogram (24) · aspect ratio · edge density · body proportions · symmetry · HOG stats (3) · LBP (16).",
         font=FONT_MONO, size=11, color=INK_SOFT)

# Slide 6: Live Demo
s = slide_base()
section_header(s, "05 / DEMO", "Live Demo")
footer(s); page_number(s, 6, TOTAL)

add_text(s, Inches(0.8), Inches(2.2), Inches(12), Inches(0.7),
         "https://hanfuzhao781-carspec-ai.hf.space",
         font=FONT_MONO, size=20, color=ACCENT, bold=True)

demos = [
    ("Upload", "Vehicle exterior photo"),
    ("Predict", "3 attributes in one pass"),
    ("Compare", "Classical vs Deep side-by-side"),
    ("Explain", "50-D feature narrative"),
]
for i, (title, desc) in enumerate(demos):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x, Inches(3.2), Inches(2.9), Inches(2.2), PAPER, line=INK)
    add_text(s, x + Inches(0.2), Inches(3.35), Inches(2.5), Inches(0.4),
             f"0{i+1}", font=FONT_MONO, size=12, color=ACCENT, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.5),
             title, font=FONT_DISPLAY, size=20, bold=True)
    add_text(s, x + Inches(0.2), Inches(4.3), Inches(2.5), Inches(0.8),
             desc, font=FONT_DISPLAY, size=12, color=INK_SOFT)

add_rect(s, Inches(0.8), Inches(6.0), Inches(12), Inches(0.8), INK)
add_text(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.8),
         'Health check:  GET /health  →  {"status":"ok","models_loaded":7}',
         font=FONT_MONO, size=13, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)

# Slide 7: Results
s = slide_base()
section_header(s, "06 / RESULTS", "Results")
footer(s); page_number(s, 7, TOTAL)

headers = ["Model", "car_type", "door_count", "seat_count"]
data = [
    ("Naive (majority)", "0.229", "0.577", "0.576"),
    ("Classical (RF)", "0.407", "0.611", "0.609"),
    ("Deep (MobileNetV2)", "0.771", "0.819", "0.860"),
]
col_x = [Inches(0.8), Inches(5.0), Inches(8.0), Inches(11.0)]
col_w = [Inches(4.0), Inches(2.9), Inches(2.9), Inches(1.7)]
y = Inches(2.2)
rh = Inches(0.65)
for ci, h in enumerate(headers):
    add_rect(s, col_x[ci], y, col_w[ci], rh, INK, line=INK)
    add_text(s, col_x[ci] + Inches(0.2), y, col_w[ci] - Inches(0.4), rh,
             h, font=FONT_MONO, size=12, bold=True, color=PAPER, anchor=MSO_ANCHOR.MIDDLE)
y += rh
for ri, row in enumerate(data):
    is_best = ri == 2
    fill = ACCENT if is_best else (PAPER if ri % 2 == 0 else RGBColor(0xF0, 0xED, 0xE7))
    text_color = PAPER if is_best else INK
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], y, col_w[ci], rh, fill, line=RULE_SOFT)
        add_text(s, col_x[ci] + Inches(0.2), y, col_w[ci] - Inches(0.4), rh,
                 cell, font=FONT_MONO if ci > 0 else FONT_DISPLAY,
                 size=12 if ci > 0 else 14, bold=is_best or ci == 0,
                 color=text_color, anchor=MSO_ANCHOR.MIDDLE)
    y += rh

add_text(s, Inches(0.8), Inches(5.0), Inches(12), Inches(0.5),
         "Top-5 accuracy = 1.000 across all three tasks for the deep model.",
         font=FONT_MONO, size=12, color=ACCENT, bold=True)

add_rect(s, Inches(0.8), Inches(5.7), Pt(4), Inches(1.0), ACCENT)
add_text(s, Inches(1.1), Inches(5.7), Inches(11.5), Inches(1.0),
         "Evaluated on 969 held-out real photos (80/20 split of 4,869 Bing-crawled images). Deep model correctly predicts 5/5 demo samples with confidence 0.97-0.98. Scaling 100 -> 4,869 lifted car_type accuracy 0.55 -> 0.77 (+22 pts).",
         font=FONT_DISPLAY, size=14, italic=True, color=INK_SOFT)

# Slide 8: Insights
s = slide_base()
section_header(s, "07 / INSIGHTS", "Findings")
footer(s); page_number(s, 8, TOTAL)

insights = [
    ("01", "Multi-task lifts door/seat", "0.82/0.86 vs 0.77 car_type. Shared backbone picks up car_type<->door<->seat correlations."),
    ("02", "Deep beats classical by 36 pts", "0.77 vs 0.41 on car_type. Learned features beat handcrafted on real photos."),
    ("03", "Classical doubles naive", "0.41 vs 0.23. 50-D handcrafted features carry real discriminative signal."),
    ("04", "Robustness under noise", "Mean accuracy 0.543 across 3 gaussian noise severity levels; classical barely moves."),
    ("05", "Handcrafted features read out as sentences", "\"dominant color is red\", \"aspect ratio 1.2 -> leans SUV\"."),
]
y = Inches(2.1)
for num, title, desc in insights:
    add_rect(s, Inches(0.8), y, Inches(12), Pt(1), RULE_SOFT)
    add_text(s, Inches(0.8), y + Inches(0.15), Inches(0.6), Inches(0.4),
             num, font=FONT_MONO, size=14, color=ACCENT, bold=True)
    add_text(s, Inches(1.6), y + Inches(0.15), Inches(4.5), Inches(0.4),
             title, font=FONT_DISPLAY, size=15, bold=True)
    add_text(s, Inches(6.3), y + Inches(0.15), Inches(6.5), Inches(0.7),
             desc, font=FONT_DISPLAY, size=12, color=INK_SOFT)
    y += Inches(0.85)
add_rect(s, Inches(0.8), y, Inches(12), Pt(1), RULE_SOFT)

# Slide 9: Engineering
s = slide_base()
section_header(s, "08 / ENGINEERING", "Engineering")
footer(s); page_number(s, 9, TOTAL)

items = [
    ("Code", "14 modular scripts (crawl -> clean -> train -> eval -> deploy)"),
    ("Models", "6 .pkl + 1 .pt on HF Hub (auto-downloaded at runtime)"),
    ("Deployment", "HF Space · gunicorn · Python 3.11 · Dockerfile"),
    ("CI", "keep-alive cron every 6h"),
    ("Git", "24 PRs · branch-based workflow · PR template"),
    ("Reports", "TECHNICAL_REPORT.md (15 sections) · GRADING.md · README.md"),
]
y = Inches(2.2)
for label, desc in items:
    add_rect(s, Inches(0.8), y, Inches(0.05), Inches(0.5), ACCENT)
    add_text(s, Inches(1.1), y, Inches(3.0), Inches(0.5),
             label, font=FONT_MONO, size=12, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.3), y, Inches(8.5), Inches(0.5),
             desc, font=FONT_DISPLAY, size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)

# Slide 10: Future + Thank You
s = slide_base(paper=False)
add_rect(s, 0, 0, SW, SH, INK)
add_mono_label(s, Inches(0.8), Inches(0.6), "09 / FUTURE", color=ACCENT)
add_text(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
         "If I had another semester...", font=FONT_DISPLAY, size=36, bold=True,
         color=PAPER, italic=True)
add_rect(s, Inches(0.8), Inches(1.7), Inches(2), Pt(3), ACCENT)

futures = [
    "Train on full 136k CompCars images with GPU",
    "Grad-CAM heatmaps for CNN attention visualization",
    "SE-Block / CBAM attention modules",
    "Multi-view fusion using CompCars view annotations",
    "TensorRT / ONNX for real-time inference",
]
y = Inches(2.2)
for i, f in enumerate(futures):
    add_text(s, Inches(0.8), y, Inches(0.6), Inches(0.4),
             f"0{i+1}", font=FONT_MONO, size=14, color=ACCENT, bold=True)
    add_text(s, Inches(1.6), y, Inches(11), Inches(0.5),
             f, font=FONT_DISPLAY, size=16, color=PAPER)
    y += Inches(0.55)

# thank you
add_rect(s, Inches(0.8), Inches(5.6), Inches(12), Pt(2), ACCENT)
add_text(s, Inches(0.8), Inches(5.9), Inches(12), Inches(1.0),
         "Thank you. Questions?", font=FONT_DISPLAY, size=48, bold=True, color=PAPER)
add_text(s, Inches(0.8), Inches(6.9), Inches(12), Inches(0.4),
         "hanfuzhao781-carspec-ai.hf.space  ·  github.com/hanfuzhao/carspec-ai",
         font=FONT_MONO, size=12, color=ACCENT)
page_number(s, 10, TOTAL)

# Save
out = Path(__file__).resolve().parent.parent / "PITCH_SLIDES.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
